import pptx
from pptx import Presentation
import pandas as pd

def create_presentation(df, customer, quarter):
  """Creates a presentation for the specified customer and solutions.

  Args:
    df: The DataFrame containing the data.
    customer: The name of the customer.
    quarter: The quarter.

  Returns:
    The presentation object.
  """

  # Create a presentation object
  prs = Presentation()


  # Create a slide for the quarter
  slide = prs.slides.add_slide()

  # Get the first row of the data for the 
  row = df.loc[df[quarter] == True].iloc[0]

  # Create a dictionary mapping placeholder names to data
  data_mapping = {
          'Title 1': f"{customer} generates {row[f'{quarter} - Revenue Generation']} in {quarter}",
          'Text Placeholder 6': row["Customer's Goal - Objective"],
          'Text Placeholder 4': row[f'{quarter} - Cost Savings'],
          'Text Placeholder 3': row[f'{quarter} - Revenue Generation']
      }

  # Add a textbox to the slide for each column of data
  for i, shape in enumerate(slide.placeholders):
    if shape.has_text_frame and shape.name in data_mapping:
      shape.text = str(data_mapping[shape.name])

  return prs

if __name__ == "__main__":
  # Load the data
  df = pd.read_excel("excel.xlsx")

  # Get the customer and quarter
  customer = "Client"
  quarter = "Annee"

  # Create the presentation
  prs = create_presentation(df, customer, quarter)

  # Save the presentation
  prs.save(f"{customer}.pptx")


# the function will take in revenue, expenses and profit and loss arguemnts from a pandas df and will create graphs and tables for the presentation
# this will be for the clients with different data
def create_presentation2(df, year, revenue, expenses, profit_loss):
  """Creates a presentation for the specified customer.

  Args:
    df: The DataFrame containing the data.
    customer: The name of the customer.
    year: The year.
    revenue: The revenue column name.
    expenses: The expenses column name.
    profit_loss: The profit/loss column name.

  Returns:
    The presentation object.
  """

  # Create a presentation object
  prs = Presentation()

  # Create a slide for the profit and loss statement
  slide = prs.slides.add_slide()

  # Get the data for the profit and loss statement
  data = df[[revenue, expenses, profit_loss]].sum()

  # Create a table for the profit and loss statement
  table = slide.shapes.add_table(4, 2, 0, 0, 0, 0).table

  # Set the column headers
  table.cell(0, 0).text = "Description"
  table.cell(0, 1).text = "Montant"

  # Set the row labels and data
  table.cell(1, 0).text = "Revenue"
  table.cell(1, 1).text = str(data[revenue])

  table.cell(2, 0).text = "Depenses"
  table.cell(2, 1).text = str(data[expenses])

  table.cell(3, 0).text = "Profits/Pertes"
  table.cell(3, 1).text = str(data[profit_loss])

  # Create a slide for the revenue and expenses graph
  slide = prs.slides.add_slide()

  # Create a pie chart for the revenue and expenses
  chart_data = pptx.chart.data.ChartData()
  chart_data.categories = ["Revenues", "Depenses"]
  chart_data.add_series("Amount", (data[revenue], data[expenses]))
  
  x, y, cx, cy = 0, 0, 5000000, 4000000
  chart = slide.shapes.add_chart(
      pptx.enum.chart.XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
  ).chart

  chart.has_legend = True

  # Set the title of the chart
  chart.has_title = True
  chart.chart_title.text_frame.text = "Revenues et Depenses"

  return prs